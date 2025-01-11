import streamlit as st
import yfinance as yf
import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import io

def fetch_stock_data(ticker, period='1y'):
    """Fetch stock data using yfinance."""
    stock = yf.Ticker(ticker)
    return stock.history(period=period)

def create_visualizations(data, ticker):
    """Create and return stock data visualizations."""
    plots = {}

    # Interactive Line Chart for Closing Prices
    fig_close = go.Figure()
    fig_close.add_trace(go.Scatter(x=data.index, y=data['Close'], mode='lines', name='Close Price'))
    fig_close.update_layout(title=f'{ticker} Closing Prices', xaxis_title='Date', yaxis_title='Price')
    plots['Closing Prices'] = fig_close

    # Interactive Bar Chart for Volume
    fig_volume = go.Figure()
    fig_volume.add_trace(go.Bar(x=data.index, y=data['Volume'], name='Volume'))
    fig_volume.update_layout(title=f'{ticker} Volume', xaxis_title='Date', yaxis_title='Volume')
    plots['Volume'] = fig_volume

    # Candlestick Chart
    fig_candlestick = go.Figure(data=[go.Candlestick(
        x=data.index,
        open=data['Open'], high=data['High'],
        low=data['Low'], close=data['Close']
    )])
    fig_candlestick.update_layout(title=f'{ticker} Candlestick Chart', xaxis_title='Date', yaxis_title='Price')
    plots['Candlestick Chart'] = fig_candlestick

    return plots

def create_powerpoint(plots, ticker):
    """Create a PowerPoint presentation with stock data visualizations."""
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Stock Analysis for {ticker}"
    subtitle.text = "Generated using yfinance and Streamlit"

    # Add slides for each plot
    for title, plot in plots.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title

        # Add the plot image
        img_path = io.BytesIO()
        plot.write_image(img_path, format='png')
        img_path.seek(0)
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    # Save PowerPoint to BytesIO
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# Streamlit App
def main():
    st.title("Stock Analysis and PowerPoint Generator")

    st.sidebar.header("User Input")
    stocks = st.sidebar.text_input("Enter stock tickers (comma-separated):", "AAPL, MSFT")
    period = st.sidebar.selectbox("Select the period for data:", ['1mo', '3mo', '6mo', '1y', '2y', '5y'], index=3)

    if st.sidebar.button("Generate Report"):
        tickers = [ticker.strip() for ticker in stocks.split(",")]
        ppt_files = {}

        for ticker in tickers:
            st.write(f"Fetching data for {ticker}...")
            try:
                data = fetch_stock_data(ticker, period)
                plots = create_visualizations(data, ticker)

                ppt_files[ticker] = create_powerpoint(plots, ticker)
                st.success(f"Data and visualizations for {ticker} prepared!")

                # Show visualizations
                for title, plot in plots.items():
                    st.plotly_chart(plot)

            except Exception as e:
                st.error(f"Error fetching data for {ticker}: {e}")

        if ppt_files:
            for ticker, ppt_file in ppt_files.items():
                st.download_button(
                    label=f"Download {ticker} Report", 
                    data=ppt_file, 
                    file_name=f"{ticker}_report.pptx", 
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

if __name__ == "__main__":
    main()
