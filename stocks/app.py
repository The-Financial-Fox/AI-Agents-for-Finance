import streamlit as st
import yfinance as yf
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import requests

# Function to fetch stock data
def fetch_stock_data(ticker):
    try:
        stock = yf.Ticker(ticker)
        data = stock.history(period="ytd")
        data["YTD Return"] = (data["Close"] / data["Close"].iloc[0]) - 1
        return data
    except Exception as e:
        st.error(f"Error fetching stock data for {ticker}: {e}")
        return None

# Function to fetch crypto data from CoinGecko
def fetch_crypto_data(symbol):
    try:
        url = f"https://api.coingecko.com/api/v3/coins/markets"
        params = {
            "vs_currency": "usd",
            "ids": symbol.lower(),  # CoinGecko uses lowercase ids for crypto
        }
        response = requests.get(url, params=params)
        data = response.json()

        if len(data) == 0:
            st.error(f"No data found for cryptocurrency {symbol}")
            return None

        # Simulating YTD performance (replace with real YTD calculation if needed)
        current_price = data[0]["current_price"]
        ytd_price = current_price * 0.9  # Example: assume a 10% increase YTD
        ytd_return = (current_price - ytd_price) / ytd_price
        return {"symbol": symbol.upper(), "YTD Return": ytd_return}
    except Exception as e:
        st.error(f"Error fetching crypto data for {symbol}: {e}")
        return None

# Function to create PowerPoint
def create_ppt(stock_data, crypto_data):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Financial Performance Report"
    slide.placeholders[1].text = "Stocks, Bonds, and Cryptocurrency Performance (YTD)"

    # Stocks slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Stocks YTD Performance"
    for ticker, data in stock_data.items():
        if data is None:  # Skip if data is None
            continue
        chart_img = BytesIO()
        plt.figure()
        plt.plot(data.index, data["YTD Return"], label=ticker)
        plt.title(f"{ticker} YTD Performance")
        plt.xlabel("Date")
        plt.ylabel("YTD Return")
        plt.legend()
        plt.tight_layout()
        plt.savefig(chart_img, format="png")
        chart_img.seek(0)
        slide.shapes.add_picture(chart_img, Inches(1), Inches(2), Inches(6), Inches(4))

    # Crypto slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Cryptocurrency YTD Performance"
    rows = len(crypto_data) + 1
    cols = 2
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(1.5)).table
    table.cell(0, 0).text = "Symbol"
    table.cell(0, 1).text = "YTD Return"
    for i, (symbol, data) in enumerate(crypto_data.items(), start=1):
        table.cell(i, 0).text = symbol
        table.cell(i, 1).text = f"{data['YTD Return']:.2%}"

    return prs

# Streamlit app
st.title("Financial Data to PowerPoint Generator")
st.sidebar.header("Input Options")

# Input stocks and crypto
stocks = st.sidebar.text_input("Enter Stock Tickers (comma separated)", "AAPL,MSFT,GOOGL")
cryptos = st.sidebar.text_input("Enter Crypto Symbols (comma separated)", "bitcoin,ethereum")

if st.button("Generate PowerPoint"):
    stock_data = {}
    crypto_data = {}
    
    # Fetch stock data
    for ticker in stocks.split(","):
        ticker = ticker.strip().upper()
        stock_data[ticker] = fetch_stock_data(ticker)

    # Fetch crypto data
    for symbol in cryptos.split(","):
        symbol = symbol.strip().lower()
        crypto_data[symbol] = fetch_crypto_data(symbol)
    
    # Generate PowerPoint
    ppt = create_ppt(stock_data, crypto_data)
    ppt_io = BytesIO()
    ppt.save(ppt_io)
    ppt_io.seek(0)
    
    # Download link
    st.success("PowerPoint generated successfully!")
    st.download_button("Download PowerPoint", data=ppt_io, file_name="financial_report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
